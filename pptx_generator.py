from __future__ import annotations

from datetime import datetime
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from chart_generator import (
    chart_baibu_compare,
    chart_link_roi_rank,
    chart_product_profit_contrib,
    chart_product_sales_top10,
    chart_q2_progress,
)
from report_pack_parser import parse_report_pack


COLORS = {
    "dark": RGBColor(0x23, 0x4B, 0x2C),
    "light": RGBColor(0xEE, 0xF6, 0xE7),
    "accent": RGBColor(0x00, 0xAC, 0xF8),
    "gray": RGBColor(0x6E, 0x7F, 0x6A),
}


def _title(slide, text: str, page: int):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(11.5), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS["dark"]

    foot = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.5), Inches(0.4))
    fp = foot.text_frame.paragraphs[0]
    fp.text = f"艾兰得经营分析系统 ｜ {page}/10"
    fp.font.size = Pt(11)
    fp.font.color.rgb = COLORS["gray"]


def _metrics_box(slide, rows: list[str], left=0.6, top=1.2, width=5.8, height=4.8):
    shp = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = COLORS["light"]
    shp.line.color.rgb = COLORS["dark"]
    tf = shp.text_frame
    tf.clear()
    for i, line in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)


def _to_float(value) -> float | None:
    if isinstance(value, (int, float)):
        return float(value)
    if isinstance(value, str):
        txt = value.replace(",", "").replace("%", "").replace("¥", "").strip()
        if not txt:
            return None
        try:
            return float(txt)
        except ValueError:
            return None
    return None


def _fmt_amount(value) -> str:
    num = _to_float(value)
    return f"¥{num:,.0f}" if num is not None else "N/A"


def _fmt_percent(value) -> str:
    num = _to_float(value)
    return f"{num:.2f}%" if num is not None else "N/A"


def _fmt_roi(value) -> str:
    num = _to_float(value)
    return f"{num:.2f}" if num is not None else "N/A"


def _add_data_table(slide, df, cols: list[str], left: float, top: float, width: float, height: float):
    show_cols = [c for c in cols if c in df.columns]
    if not show_cols:
        return
    preview = df[show_cols].head(6).fillna("-")
    table = slide.shapes.add_table(len(preview) + 1, len(show_cols), Inches(left), Inches(top), Inches(width), Inches(height)).table
    for i, c in enumerate(show_cols):
        cell = table.cell(0, i)
        cell.text = str(c)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)
    for r_idx, (_, row) in enumerate(preview.iterrows(), start=1):
        for c_idx, col in enumerate(show_cols):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(row[col])
            cell.text_frame.paragraphs[0].font.size = Pt(10)


def build_pptx_from_report_pack(report_pack: dict) -> bytes:
    parsed = parse_report_pack(report_pack)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    overview, q2 = parsed["overview_metrics"], parsed["q2_kpi"]
    if parsed["product_table"].empty or parsed["link_table"].empty or parsed["baibu_table"].empty:
        raise ValueError("产品、链接、百补页面必须包含真实数据（含表格行），请从经营系统导出完整JSON数据包后重试。")

    # 1 cover
    s = prs.slides.add_slide(blank)
    _title(s, "艾兰得拼多多经营分析报告", 1)
    _metrics_box(s, ["经营总览 / Q2考核 / 推广效率 / 异常清单", f"生成时间：{datetime.now():%Y-%m-%d %H:%M}", "渠道：拼多多"], 0.8, 2.1, 11.5, 2.2)

    # 2 overview
    s = prs.slides.add_slide(blank)
    _title(s, "经营总览", 2)
    rows = [
        f"商家实收: {_fmt_amount(overview.get('商家实收'))}",
        f"用户实付: {_fmt_amount(overview.get('用户实付'))}",
        f"有效订单数: {overview.get('有效订单数', 'N/A')}",
        f"订单侧估算毛利: {_fmt_amount(overview.get('订单侧估算毛利'))}",
        f"店铺总盘推广费（现金口径）: {_fmt_amount(overview.get('店铺总盘推广费（现金口径）'))}",
        f"店铺整体实际ROI: {_fmt_roi(overview.get('店铺整体实际ROI'))}",
    ]
    _metrics_box(s, rows)

    # 3 q2
    s = prs.slides.add_slide(blank)
    _title(s, "Q2考核达成率", 3)
    q2_rows = [
        f"当前销售额: {_fmt_amount(q2.get('当前销售额'))}",
        f"Q2销售目标: {_fmt_amount(q2.get('Q2销售目标'))}",
        f"销售达成率: {_fmt_percent(q2.get('销售达成率'))}",
        f"当前实际ROI: {_fmt_roi(q2.get('当前实际ROI'))}",
        f"Q2 ROI目标: {_fmt_roi(q2.get('Q2 ROI目标'))}",
        f"ROI达成率: {_fmt_percent(q2.get('ROI达成率'))}",
        f"综合达成率: {_fmt_percent(q2.get('综合达成率'))}",
        f"奖金风险等级: {q2.get('奖金风险等级', 'N/A')}",
    ]
    _metrics_box(s, q2_rows, 0.6, 1.1, 5.8, 5.5)
    img = chart_q2_progress(q2)
    if img:
        s.shapes.add_picture(BytesIO(img), Inches(6.7), Inches(1.4), width=Inches(6.0))
    target = q2.get("Q2销售目标", 0)
    if float(str(target).replace('%','') or 0) == 0:
        warn = s.shapes.add_textbox(Inches(6.7), Inches(5.9), Inches(6), Inches(0.8))
        warn.text_frame.text = "当前数据包未包含有效Q2销售目标，请回经营系统填写后重新导出。"

    # 4 gap
    s = prs.slides.add_slide(blank)
    _title(s, "达标缺口测算", 4)
    if float(str(q2.get("Q2销售目标", 0)).replace('%','') or 0) == 0:
        _metrics_box(s, ["当前数据包未包含有效 Q2销售目标，请回经营系统填写 Q2销售目标后重新导出。"], 0.8, 2.0, 11.4, 1.4)
    else:
        gap_rows = [f"{k}: {q2.get(k, 'N/A')}" for k in ["70%销售安全线", "距离70%安全线还差", "距离100%销售目标还差", "达到100%目标每日需完成销售额", "距离目标ROI还差"]]
        _metrics_box(s, gap_rows)

    # 5 product sales
    s = prs.slides.add_slide(blank)
    _title(s, "产品表现 Top 10", 5)
    img = chart_product_sales_top10(parsed["product_table"])
    if img:
        s.shapes.add_picture(BytesIO(img), Inches(0.7), Inches(1.1), width=Inches(8.8))
    _add_data_table(s, parsed["product_table"], ["标准产品名称", "商家实收", "扣推广后贡献毛利", "实际ROI"], 0.7, 5.15, 8.8, 1.9)
    _metrics_box(s, ["经营结论：聚焦头部产品保增长，尾部SKU控投放。"], 9.7, 1.5, 3.0, 2.0)

    # 6 profit
    s = prs.slides.add_slide(blank)
    _title(s, "产品利润贡献分析", 6)
    img = chart_product_profit_contrib(parsed["product_table"])
    if img:
        s.shapes.add_picture(BytesIO(img), Inches(0.8), Inches(1.2), width=Inches(11.8))

    # 7 link
    s = prs.slides.add_slide(blank)
    _title(s, "链接表现分析", 7)
    img = chart_link_roi_rank(parsed["link_table"])
    if img:
        s.shapes.add_picture(BytesIO(img), Inches(0.6), Inches(1.1), width=Inches(8.0))
    _add_data_table(s, parsed["link_table"], ["链接标题", "商家实收", "推广费", "实际ROI"], 8.7, 1.2, 4.2, 4.8)
    _metrics_box(s, ["经营结论：优先追加高ROI链接预算，暂停低ROI且负毛利链接。"], 8.7, 6.1, 4.2, 1.0)

    # 8 baibu
    s = prs.slides.add_slide(blank)
    _title(s, "百补 vs 日常", 8)
    img = chart_baibu_compare(parsed["baibu_table"])
    if img:
        s.shapes.add_picture(BytesIO(img), Inches(0.8), Inches(1.2), width=Inches(8.5))
    _add_data_table(s, parsed["baibu_table"], ["类型", "商家实收", "推广费", "扣推广后贡献毛利", "实际ROI"], 8.9, 1.2, 3.9, 4.8)
    _metrics_box(s, ["经营结论：判断百补偏规模还是利润，识别亏损风险更高侧。"], 9.2, 6.1, 3.3, 1.0)

    # 9 abnormal
    s = prs.slides.add_slide(blank)
    _title(s, "经营异常", 9)
    abn = parsed["abnormal_table"]
    lines = ["异常清单："]
    if not abn.empty:
        lines += [str(r) for r in abn.head(8).to_dict(orient="records")]
    else:
        lines += ["未提供异常表，已建议优先排查负毛利产品/链接。"]
    _metrics_box(s, lines, 0.6, 1.2, 12.0, 5.8)

    # 10 advice
    s = prs.slides.add_slide(blank)
    _title(s, "下阶段经营动作建议", 10)
    base = ["先突破70%销售安全线", "控制低ROI推广费", "优化百补亏损链接", "放大利润规格", "提升高毛利产品曝光", "每周复盘产品/链接 ROI"]
    extra = [i.strip() for i in str(parsed["advice_text"]).replace("；", "\n").splitlines() if i.strip()]
    actions = (extra + base)[:6]
    _metrics_box(s, [f"- {a}" for a in actions], 0.8, 1.4, 11.5, 4.8)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()
